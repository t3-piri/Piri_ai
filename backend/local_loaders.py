import re

from docx import Document
from docx.oxml.ns import qn
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

HEADING_RE = re.compile(r"^\s*(\d+(\.\d+){0,3}\.?)\s+\S")


def _iter_block_items(document):
    body = document.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, document)
        elif child.tag == qn("w:tbl"):
            yield Table(child, document)


def load_docx(path):
    """Returns list of (text, locator) blocks in document order."""
    doc = Document(path)
    blocks = []
    heading = None
    for item in _iter_block_items(doc):
        if isinstance(item, Paragraph):
            text = item.text.strip()
            if not text:
                continue
            style = (item.style.name or "") if item.style else ""
            is_heading_style = "Heading" in style or "Title" in style
            is_numbered_clause = bool(HEADING_RE.match(text)) and len(text) < 120
            if is_heading_style or is_numbered_clause:
                heading = text[:100]
            blocks.append((text, heading))
        elif isinstance(item, Table):
            rows = []
            for row in item.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                blocks.append(("\n".join(rows), heading))
    return blocks


def load_pdf(path):
    """Returns list of (text, locator) blocks, one per page."""
    reader = PdfReader(path)
    blocks = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            blocks.append((text, f"Sayfa {i}"))
    return blocks


def load_pptx(path):
    """Returns list of (text, locator) blocks, one per slide.
    Separates the slide title from body content and appends speaker notes."""
    prs = Presentation(path)
    blocks = []
    for i, slide in enumerate(prs.slides, start=1):
        title = None
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None and title_shape.has_text_frame:
            title = title_shape.text_frame.text.strip() or None

        body_parts = []
        for shape in slide.shapes:
            if shape is title_shape:
                continue
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    body_parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        body_parts.append(" | ".join(cells))

        notes = None
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                notes = notes_text

        parts = []
        if title:
            parts.append(f"Başlık: {title}")
        parts.extend(body_parts)
        if notes:
            parts.append(f"Notlar: {notes}")

        if parts:
            blocks.append(("\n".join(parts), f"Slayt {i}"))
    return blocks


def _row_to_sentence(header, cells, sheet_title):
    pairs = [f"{h.strip()}: {c.strip()}" for h, c in zip(header, cells) if h and c]
    if not pairs:
        return ""
    return f"{sheet_title} tablosunda şu kayıt var -> " + "; ".join(pairs) + "."


def load_xlsx(path):
    """Returns list of (text, locator) blocks. Each row becomes a readable
    sentence built from its header labels, instead of a raw '|' separated row."""
    wb = load_workbook(path, data_only=True, read_only=True)
    blocks = []
    for sheet in wb.worksheets:
        header = None
        for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            cells = ["" if c is None else str(c).strip() for c in row]
            if not any(cells):
                continue
            if header is None:
                header = cells
                continue
            sentence = _row_to_sentence(header, cells, sheet.title)
            if sentence:
                blocks.append((sentence, f"{sheet.title}, satır {row_idx}"))
    return blocks


LOADERS = {
    ".docx": load_docx,
    ".pdf": load_pdf,
    ".pptx": load_pptx,
    ".xlsx": load_xlsx,
}


def load_file(path):
    ext = path.suffix.lower()
    loader = LOADERS.get(ext)
    if loader is None:
        return []
    try:
        return loader(str(path))
    except Exception as e:
        print(f"UYARI: {path} okunamadi: {e}")
        return []
